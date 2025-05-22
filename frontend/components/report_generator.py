from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_text_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    textbox = slide.placeholders[1]
    textbox.text = content

def add_table_slide(prs, title, df):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)).table

    # 헤더
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # 내용
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            cell.text = str(df.iloc[i, j])
            cell.text_frame.paragraphs[0].font.size = Pt(12)

def add_positioning_slide(prs, comp_id, df, overall_winner, overall_judgement):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Technology Positioning vs {comp_id}"

    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.0), Inches(9), Inches(4)).table
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            cell.text = str(df.iloc[i, j])
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    # 결론 텍스트 박스 추가
    left = Inches(0.5)
    top = Inches(5.2)
    width = Inches(9)
    height = Inches(1.2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p1 = tf.add_paragraph()
    p1.text = f"🏁 Overall Winner: {overall_winner}"
    p1.font.size = Pt(14)
    p2 = tf.add_paragraph()
    p2.text = f"🧠 Reason: {overall_judgement}"
    p2.font.size = Pt(12)

def add_implementation_slide(prs, comp_id, df, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Implementation Difference vs {comp_id}"

    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.0), Inches(9), Inches(4)).table
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            cell.text = str(df.iloc[i, j])
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    left = Inches(0.5)
    top = Inches(5.2)
    width = Inches(9)
    height = Inches(1.2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = f"🧠 Summary of Technical Difference: {summary}"
    p.font.size = Pt(12)

def generate_ppt_from_result(strategy_output, pos_result, imp_diff_result, output_path="report.pptx", out_stream=None):
    prs = Presentation()

    # 1. Title slide
    add_title_slide(prs, "Patent Comparison Report", "3-Layer Based Patent Analysis")

    # 2. Strategic Recommendation
    add_text_slide(prs, "Strategic Recommendation", strategy_output["our_overall_strategy"])

    # 3. Strategic Summary Table
    df_strategy = pd.DataFrame(strategy_output["strategy_table"])
    table_df = df_strategy[["patent_id", "tech_summary", "tech_similarity", "technical_value", "strategic_direction"]]
    add_table_slide(prs, "Patent Strategy Summary Table", table_df)

    # 4. Technology Positioning Result
    for df in pos_result:
        comp_id = df.attrs.get("competitor_id", "Unknown")
        overall_winner = df.attrs.get("overall_winner", "N/A")
        overall_judgement = df.attrs.get("overall_judgement", "N/A")
        add_positioning_slide(prs, comp_id, df, overall_winner, overall_judgement)

    # 5. Implementation Difference Result
    for df in imp_diff_result:
        comp_id = df.attrs.get("competitor_id", "Unknown")
        summary = df.attrs.get("overall_diff_summary", "N/A")
        add_implementation_slide(prs, comp_id, df, summary)

    if out_stream is not None:
        prs.save(out_stream)
    else:
        prs.save(output_path)
    return output_path or out_stream
